#-*- coding: utf-8 -*-
##＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊
import tkinter as tk
import datetime

class Application(tk.Frame):
    def __init__(self, master=None):
        super().__init__(master)
        self.master = master
        self.pack()

        self.create_widgets()
        self.scan_completed = False
        self.nsite = 0
        self.nslot = 0
        self.rawdatas = []
        self.rawdatas_for_table = []
        #self.scan_data_files()
        self.df2 = []
        self.n_report_sets_per_page = 1
    
    def create_widgets(self):
        self.lbl_date = tk.Label(self)
        self.lbl_date["text"]="日期(Date): [YYYY/MM/DD]"
        self.lbl_date.grid(row=0, column=0)
        
        self.txt_report_date = tk.Entry(self)
        ## 取得現在時間
        xtime=datetime.datetime.now()                                                
        ## 建立一個 tk string變數物件
        report_date = tk.StringVar()                                                 
        ## 設定預設字串
        report_date.set("{}/{:02d}/{:02d}".format(xtime.year,xtime.month,xtime.day)) 
        ## entry物件將預設字串賦值
        self.txt_report_date["textvariable"]=report_date
        self.txt_report_date.grid(row=0, column=1)
        
        self.lbl_report_editor = tk.Label(self)
        self.lbl_report_editor["text"]="作者(Editor): "
        self.lbl_report_editor.grid(row=1,column=0)
        self.txt_report_editor = tk.Entry(self)
        self.txt_report_editor.grid(row=1,column=1)
        
        self.lbl_report_to_customer = tk.Label(self)
        self.lbl_report_to_customer["text"]="客戶名稱: "
        self.lbl_report_to_customer.grid(row=2,column=0)
        self.txt_report_to_customer = tk.Entry(self)
        self.txt_report_to_customer.grid(row=2,column=1)
        
        self.lbl_product_id = tk.Label(self)
        self.lbl_product_id["text"] = "產品型號(Product id): "
        self.lbl_product_id.grid(row=3,column=0)
        self.txt_product_id = tk.Entry(self)
        self.txt_product_id.grid(row=3,column=1)
        
        self.lbl_lot_id = tk.Label(self)
        self.lbl_lot_id["text"] = "產品批號(lot_id): "
        self.lbl_lot_id.grid(row=4, column=0)
        self.txt_lot_id = tk.Entry(self)
        self.txt_lot_id.grid(row=4, column=1)
        
        self.lbl_stage_id = tk.Label(self)
        self.lbl_stage_id["text"] = "產品階段(stage_id): [What process?]"
        self.lbl_stage_id.grid(row=5, column=0)
        self.txt_stage_id = tk.Entry(self)
        self.txt_stage_id.grid(row=5, column=1)
        ## ===== 2021/6/16 增加輸出版面的參數: =====
        self.lbl_report_sets_per_page = tk.Label(self)
        self.lbl_report_sets_per_page["text"] = "每頁放多少片數:"
        self.lbl_report_sets_per_page.grid(row=6, column=0)
        self.txt_report_sets_per_page = tk.Entry(self)
        self.txt_report_sets_per_page.grid(row=6, column=1)
        self.txt_report_sets_per_page.insert(0,"1")
        ## ===== 2022/1/4 增加直接導入檔名量測數值 =====
        ## 增加直接導入檔名量測數值, 預設為 0(走OCR), 使用者勾選要不要導入 1(直接取value處理)
        self.int_direct_handdata = tk.IntVar()
        self.chk_handdata = tk.Checkbutton(self,onvalue=1,offvalue=0,variable=self.int_direct_handdata,command=self.set_handdata)
        self.chk_handdata["text"]="Directly input measurement data in filename?"
        self.chk_handdata.grid(row=7,column=0)
        print("[Application] checkbox=",self.int_direct_handdata.get())

        self.btn_scan_data_files = tk.Button(self)
        self.btn_scan_data_files["text"] = "Scan data files"
        self.btn_scan_data_files["command"] = self.scan_data_files
        self.btn_scan_data_files.grid(row=8,column=0)
        
        self.btn_run_ocr = tk.Button(self)
        self.btn_run_ocr["text"] = "Run\nRecognition"
        self.btn_run_ocr["command"] = self.run_ocr
        self.btn_run_ocr.grid(row=8,column=1)
        
        self.btn_make_report = tk.Button(self)
        self.btn_make_report["text"] = "Make report"
        self.btn_make_report["command"] = self.make_report
        self.btn_make_report.grid(row=8,column=2)
        self.btn_make_report["state"] = "disabled"
        self.btn_make_report.flash()
        
        self.quit = tk.Button(self, text="QUIT", fg="red",command=self.say_bye)
        self.quit.grid(row=9, column=2)
        
    def set_handdata(self):
        print("[Application] checkbox=",self.int_direct_handdata.get())

    def run_ocr(self):
        print("[Application] ===== Run OCR recognition =====\n")
        report_date = self.get_txt_report_date_value()
        #print(type(report_date),report_date)
        report_editor = self.get_txt_report_editor_value()
        #print(type(report_date),report_editor)
        report_to_customer = self.get_txt_report_to_customer_value()
        #print(report_to_customer)
        product_id = self.get_txt_product_id_value()
        #print(product_id)
        lot_id = self.get_txt_lot_id_value()
        #print(lot_id)
        stage_id = self.get_txt_stage_id_value()
        #print(stage_id)
        
        ## 建立 OCR 影像辨識物件: => ocr
        ocr = OCR(self.rawdatas)
        
        ## =====2022/1/4 增加直接導入檔名中的量測數據.=====
        self.rawdatas_for_table = ocr.go(DirectHanddata=self.int_direct_handdata.get())
        
        self.df2 = ocr.gotable()
    
        ocr.output_datafiles()
        
        self.btn_make_report["state"] = "normal"
        self.btn_make_report.flash()
        print("[Application] ***************** Recognition completed! ********************* ")
    def say_bye(self):
        print("[Application] bye bye!")
        self.master.destroy()
        
##== define the member function "get_[object_name]_value():" ====
    def get_txt_report_date_value(self):
        return self.txt_report_date.get()
    def get_txt_report_editor_value(self):
        return self.txt_report_editor.get()
    def get_txt_report_to_customer_value(self):
        return self.txt_report_to_customer.get()
    def get_txt_product_id_value(self):
        return self.txt_product_id.get()
    def get_txt_lot_id_value(self):
        return self.txt_lot_id.get()
    def get_txt_stage_id_value(self):
        return self.txt_stage_id.get()
    def get_txt_report_sets_per_page(self):
        return int(self.txt_report_sets_per_page.get())

    def scan_data_files(self):
        ###============== Process A : 列出目錄中檔案, 並 parsing成字典 ================
        from os import listdir
        from os.path import isfile, isdir, join
        output_image_path="tn"
        orig_image_path="orig"
        files = listdir(orig_image_path)

        rawdata={}
        rawdatas=[]
        slots=[]
        sites=["C","M","E","FE","1","2","3","4","5"]  ## 2021/1/6 add for additional site purpose
        mtypes=["ERO","DIS"]                          ## 2021/1/6 add for additional site purpose
        imgtypes=["jpg","png","JPG","PNG"]
        Center=0
        Middle=1
        Edge=2
        FEdge=3
        Site1=4 ## 2021/1/6 add for additional site purpose
        Site2=5 ## 2021/1/6 add for additional site purpose
        Site3=6 ## 2021/1/6 add for additional site purpose
        Site4=7 ## 2021/1/6 add for additional site purpose
        Site5=8 ## 2021/1/6 add for additional site purpose
        for ele in files:
            if ("jpg" in ele) or ("png" in ele):
                #print(ele)
                if ele[0] == "S" or ele[0] =="s":
                    slot_num = int(ele[1:3])
                    #print(slot_num)
                    if ele[4] == 'F' or ele[4] =='f':
                        site=ele[4:6]
                        mtype=ele[7:10]
                        nsite=FEdge
                    elif ele[4] =='C' or ele[4]=='c': 
                        site=ele[4]
                        mtype=ele[6:9]
                        nsite=Center
                    elif ele[4] =='M' or ele[4] =='m':
                        site=ele[4]
                        mtype=ele[6:9]
                        nsite=Middle
                    elif ele[4] =='E' or ele[4] =='e':
                        site=ele[4]
                        mtype=ele[6:9]
                        nsite=Edge
                    elif ele[4] =='1': ## 2021/1/6 add for additional site purpose
                        site=ele[4]
                        mtype=ele[6:8]
                        nsite=Site1
                    elif ele[4] =='2': ## 2021/1/6 add for additional site purpose
                        site=ele[4]
                        mtype=ele[6:8]
                        nsite=Site2
                    elif ele[4] =='3': ## 2021/1/6 add for additional site purpose
                        site=ele[4]
                        mtype=ele[6:8]
                        nsite=Site3
                    elif ele[4] =='4': ## 2021/1/6 add for additional site purpose
                        site=ele[4]
                        mtype=ele[6:8]
                        nsite=Site4
                    elif ele[4] =='5': ## 2021/1/6 add for additional site purpose
                        site=ele[4]
                        mtype=ele[6:8]
                        nsite=Site5

                    rawdata['slot']=slot_num
                    rawdata['site']=site
                    rawdata['nsite']=nsite
                    rawdata['mtype']=mtype
                    rawdata['filename']=ele
                    #print(rawdata)
                    rawdatas.append(rawdata)
                    rawdata={}
                else:
                    print("[Application] filename {} is not the data file".format(ele))
            #rawdatas.append(rawdata)        


        rawdatas.sort(key=lambda k:k['nsite']) ## 先對對應的 C,M,E,FE做排序,
        rawdatas.sort(key=lambda k:k['slot'])  ## 再對 slot做排序
        
        # =====印出字典=====
        #print("===== \n 字典list: \n",rawdatas)
        
        # =====取得片數:=====
        # 2020/09/29 原寫法若片數有缺片中斷會有 bug. ex: S07, S08, S10, S11
        #slot_num_max = max([ele['slot'] for ele in rawdatas])
        #slot_num_min = min([ele['slot'] for ele in rawdatas])
        tmp_slot = []
        for ele in rawdatas:
            if ele['slot'] not in tmp_slot:
                tmp_slot.append(ele['slot'])
        
        print("[Application] #247: ", tmp_slot)
        nslot = len(tmp_slot)
        print("[Application] #249: ", nslot)
        
        print("[Application] ===== \n 總共有: ",nslot,"片")
        # =====取得量測位置數目:=====
        site_count_max = max([ele['nsite'] for ele in rawdatas])
        site_count_min = min([ele['nsite'] for ele in rawdatas])
        nsite = site_count_max - site_count_min +1
        
        print("[Application] ===== \n 量測位置有:",nsite,"個")
        self.rawdatas=rawdatas
        self.nslot = nslot
        self.nsite=nsite
        self.scan_completed=True
        if self.scan_completed: 
            self.btn_scan_data_files["state"]="disabled"
            self.btn_scan_data_files.flash()
        
    def get_rawdatas(self):
        return self.rawdatas
    def get_nslot(self):
        return self.nslot
    def get_nsite(self):
        return self.nsite
    
    def make_report(self):
        ## 將每頁組數放入 npages全域變數
        self.n_report_sets_per_page = self.get_txt_report_sets_per_page()
        
        ## 建立報告產生物件 REPORTGEN: => report
        report = REPORTGEN(self.get_txt_report_date_value(),
                           self.get_txt_report_editor_value(),
                           self.get_txt_report_to_customer_value(),
                           self.get_txt_product_id_value(),
                           self.get_txt_lot_id_value(),
                           self.get_txt_stage_id_value(),
                           self.rawdatas_for_table, self.nsite, self.nslot, self.df2, self.n_report_sets_per_page)
        report.go()
        report.save_report()
##＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊
##＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊
import pandas as pd
import numpy as np
import cv2
import glob
import pytesseract

class OCR():
    def __init__(self, rawdatas):
        self.rawdatas = rawdatas
        self.rawdatas_for_table=[]
        self.outputlog = []
        print("[OCR] ============OCR object============\n",self.rawdatas)
        print("[OCR] ============OCR object============")
        self.df2=[]
        
    def unsharp_mask(self, image, kernel_size=(5, 5), sigma=1.0, amount=1.0, threshold=0):
        """Return a sharpened version of the image, using an unsharp mask."""
        blurred = cv2.GaussianBlur(image, kernel_size, sigma)
        sharpened = float(amount + 1) * image - float(amount) * blurred
        sharpened = np.maximum(sharpened, np.zeros(sharpened.shape))
        sharpened = np.minimum(sharpened, 255 * np.ones(sharpened.shape))
        sharpened = sharpened.round().astype(np.uint8)
        if threshold > 0:
            low_contrast_mask = np.absolute(image - blurred) < threshold
            np.copyto(sharpened, image, where=low_contrast_mask)
        return sharpened
    
    ##### ===== image process functions =====
    ##### 1張圖 呼叫執行 1次: 
    def process_img(self, filename,output_filename,width,height):
        ## 變更圖檔size => (800,547)
        image=np.array([])
        try:
            image = cv2.imread(filename)
        except:
            print("[OCR] There is not [ {} ] be found!".format(filename))
            return image
        if width == 0 or height == 0 :
            width = 800
            height = 547
        else:
            pass
        ## interpolation : 
        ## INTER_NEAREST : 鄰近插值
        ## INTER_LINEAR : 雙線性插植
        ## INTER_AREA : 像素區域關係重採樣
        ## INTER_CUBIC : 4*4像素鄰近區域雙三次插值
        ## INTER_LANCZOS4 : 8*8像素鄰近區域 Lanczos插值

        resized_image = cv2.resize(image, (width,height), interpolation=cv2.INTER_CUBIC)
        ## 將圖檔反相:
        image_n = ~resized_image

        hsv_image = cv2.cvtColor(image_n,cv2.COLOR_BGR2HSV)

        # 定義HSV下"skyblue"顏色的上下界:
        skyblue_lo=np.array([0,46,82])
        skyblue_hi=np.array([92,255,255])
        # 建立 skyblue的 mask做選擇:
        mask_image=cv2.inRange(hsv_image,skyblue_lo,skyblue_hi)

        # 將選擇的skyblue區域, 顏色變更為 (B=255,G=0,R=0)
        image_n[mask_image>0]=(255,0,0)

        result = cv2.imwrite(output_filename,image_n)
        return result

    def extract_handdata(self,fstring):
        ## 2022/1/4 直接將手輸入的數據切出來.
        if fstring != "":                             ## 若字串不為空才進來做parsing.
            r_pos = fstring.rfind("-")                ## 從右邊開始找 第1個"-"的index.
            if fstring[r_pos-1] == "-":               ## 判斷上面找到的 "-"左邊一個是不是也是 "-",
                r_pos = r_pos-1                       ##代表這個是負號, 前一個 "-"才是分隔符號. 並更新正確的 index.
            else:
                pass
            dot_pos = fstring.rfind(".")              ## 從右邊開始找第1個 "."的index.

            str_handdata = fstring[r_pos+1 : dot_pos] ##切割出量測數值作為回傳.
        else:
            str_handdata = ""                         ##空字串
        return str_handdata

    def clean_data_num(self, str_data,SHOWDEBUGMSG=False):
        if str_data == None:
            return " "
        digit_pattern = "0123456789.-"
        number_pattern = "0123456789"

        old_len_str_data = len(str_data)

        pos = 0

        space_pos=[]
        unknown_digit=[]

        str_tmp=str_data
    ####===== 把 space 處理掉 ==============================================================================
    ####56789012345678901234567890123456789012345678901234567890

        for ch in str_tmp:
            if pos == 0:
                if ch == " ":
                    # print("left strip of space: ",str_tmp=str_tmp.strip(" "))
                    str_tmp=str_tmp.strip(" ")
                if ch in digit_pattern :
                    if SHOWDEBUGMSG : print(pos,ch)
                else:
                    unknown_digit.append(ch)
                if SHOWDEBUGMSG : print(pos,ch)
            elif pos == old_len_str_data:
                if ch == " ":
                    # print("right strip of space: ",str_tmp=str_tmp.rstrip(" "))
                    str_tmp=str_tmp.rstrip(" ")
                if ch in digit_pattern :
                    if SHOWDEBUGMSG : print(pos,ch)
                else:
                    unknown_digit.append(ch)
                if SHOWDEBUGMSG : print(pos,ch)
            elif pos > 0 and pos < old_len_str_data:
                if ch == " ":
                    space_pos.append(pos)
                    if SHOWDEBUGMSG : print(space_pos, ch)

                if ch in digit_pattern :
                    if SHOWDEBUGMSG : print(pos,ch)
                else:
                    unknown_digit.append(ch)
            pos +=1
    ####56789012345678901234567890123456789012345678901234567890
        if len(space_pos) != 0:
            if (str_tmp[space_pos[0]-1] in number_pattern ) and (str_tmp[space_pos[0]+1] in number_pattern ):
                str_tmp = str_tmp.replace(" ",".")
            else:
                str_tmp = str_tmp.replace(" ","")

    ####==================================================================================================

    ####===== 把 其他非數據字元處理掉 ========================================================================
    ####56789012345678901234567890123456789012345678901234567890
        for ud in unknown_digit:
            str_tmp = str_tmp.replace(ud,"")
    ####56789012345678901234567890123456789012345678901234567890    
    ####==================================================================================================    
        return str_tmp

    #### ===== 自動辨識數據 OCR function ======
    def recog_num(self, temp_img, target_img, roi_left, roi_top, roi_width, roi_height):
        print("[OCR] 〔辨識中......〕 ")
        ## temp_img=template image name[str]
        ## target_img=target image name[str]
        ## roi_left = 0
        ## roi_top = 240
        ## roi_width = 100
        ## roi_height = 60

        # step A : 找到 St Height字樣的位置:
        template_image = cv2.imread(temp_img,0)
        load_image = cv2.imread(target_img)
        (tW, tH) = template_image.shape[::-1]
        print("[OCR] tW,tH",(tW, tH))
        ## 開始所有 loading進來的 dishing圖做 matching 找出 loc:
        test_image= cv2.cvtColor(load_image, cv2.COLOR_BGR2GRAY)

        ## crop the test_image for "St Height" finding. [Y(240~300),X(0~100)]
        test_image=test_image[roi_top:roi_top+roi_height, roi_left:roi_left+roi_width] 

        result = cv2.matchTemplate(test_image, template_image, cv2.TM_CCOEFF_NORMED)

        print("[OCR] ",target_img,"result max = {:8.5f}".format(np.max(result)))
        threshold = 0.705
        loc = np.where(result >= threshold)

        ##loc紀載資料為 tuple e.g. (y,x)
        #print(loc)
        ## 關於 zip()的使用法: https://blog.gtwang.org/programming/python-iterate-through-multiple-lists-in-parallel/ 
        count = 0
        for point in zip(*loc[::-1]):
            #cv2.rectangle(test_image, point, (point[0] + tW, point[1] + tH), (255,0,0), 0)

            count +=1
            print("[OCR] Found {} point, (x,y)= {}".format(count,point))

            print("[OCR] ",point[1],point[0])
            ## Step B : 將 test_image按照找到的 loc point [Y,X] 切割出辨認區:
            recog_image = test_image[point[1]:point[1]+18, tW-5:tW-6+40]

            ## 做一次 USM:
            #recog_image = unsharp_mask(recog_image)

            ## recog_image 餵給 PIL
            (thresh, img) = cv2.threshold(recog_image, 128, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)

            ## 用 tesseract辨識數值
            result = pytesseract.image_to_string(img)
            ## 清理 數值
            txt = self.clean_data_num(result,SHOWDEBUGMSG=False)
            
            print("[OCR] 辨識結果=",txt)
        return txt
            #cv2.imshow('Sliced recogition image',recog_image)
            #cv2.waitKey(0)
            #cv2.imshow('Result',test_image)
            #cv2.waitKey(0)
            #cv2.destroyAllWindows()
            
    def go(self,DirectHanddata=0):
        ##### ===== Process B : 辨識數據成 dict =====
        # pytesseract
        
        orig_image_path = "orig"
        rawdatas_for_table = []
        rawdata={}
        i=0
        for ele in self.rawdatas:
            rawdata = ele.copy()
            ## 呼叫 recog_num 辨識數字: 
            print("[OCR] ocr.go():", orig_image_path, ele['filename'])
            ## 2022/1/4 增加直接導入檔名中的量測數據.
            if DirectHanddata == 1:
                result = self.extract_handdata(ele['filename'])
            else:
                #result = self.recog_num('img_data\\St_Height_cr.jpg',orig_image_path+'\\'+ele['filename'],0,240,100,60)
                result = self.recog_num('img_data/St_Height_cr.jpg',orig_image_path+'/'+ele['filename'],0,240,100,60)
            ## 將辨識的數字放進 rawdata字典中, i.e. { ..., StepH: 'result'} :
            rawdata['StepH'] = result
            ## 將一筆 rawdata 加入 rawdatas_for_table [list] 中 :
            rawdatas_for_table.append(rawdata)
            rawdata={}
            print("[OCR] ",i)
            i=i+1
        # print(rawdatas_for_table)
        self.rawdatas_for_table = rawdatas_for_table
        return self.rawdatas_for_table
        
    def gotable(self):
        ##### ===== Process C : 處理 dict 成為 table =====
        # 使用 pandas
        ##### 
        pass_slot=[] ## 記憶已處過的slot
        slot_index = 0
        tmp_dict={}
        by_site={}
        ####567890123456789012345678901234567890123456789012345678901234567890
        for ele in self.rawdatas_for_table:
            #by_site={}
            curslot = ele['slot']
            if slot_index > 0:
                if curslot == pass_slot[slot_index-1]:
                    by_site[ele['site']] = ele['StepH']
                    tmp_dict["{:02d}".format(ele['slot'])] = by_site

                else:
                    by_site={}
                    by_site[ele['site']] = ele['StepH']
                    tmp_dict["{:02d}".format(ele['slot'])] = by_site
                #print(curslot)
            else:
                by_site[ele['site']] = ele['StepH']
                tmp_dict["{:02d}".format(ele['slot'])] = by_site
                #by_site={}
                #print(curslot)
            ## 記憶前一片 slot
            pass_slot.append(curslot)
            slot_index +=1
        ####567890123456789012345678901234567890123456789012345678901234567890
        # for ele in rawdatas_for_table:
        #     print(ele)
        # for ele in tmp_dict:
        #     print(ele,tmp_dict[ele])
        ##=====產生 dataframe =====
        df1 = pd.DataFrame(tmp_dict)
        #print("df1",df1)
        ##=====將df1 table轉置
        df2=df1.T
        #print("df2",df2)
        self.df2 = df2
        return self.df2
    
    def output_datafiles(self):
        ##### ===== Process D : 處理對應的圖檔轉換並輸出到 tn\資料夾下 =====
        # opencv 
        orig_image_path = "orig"
        output_image_path = "tn"
        for ele in self.rawdatas_for_table:
            #result = self.process_img(orig_image_path + "\\" + ele['filename'],output_image_path + "\\" + ele['filename'],800,547)
            result = self.process_img(orig_image_path + "/" + ele['filename'],output_image_path + "/" + ele['filename'],800,547)
            print("[OCR] OCR.output_datafiles():","輸出={}\{}".format(output_image_path, ele['filename']), "辨識數值={}".format(ele['StepH']))
            if result : 
                #print(output_image_path + "\\" + ele['filename']," ok!") 
                self.outputlog.append({output_image_path + "/" + ele['filename'] : "ok"})
        return self.outputlog
##＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊
##＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊
from pptx import Presentation
from pptx.util import Inches
# from pptx.chart.data import ChartData, CategoryChartData, XyChartData, BubbleChartData
# from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Cm, Pt
# from config.setting import _EMUS_PER_CM
from pptx.dml.color import RGBColor
#from pythonwin.pywin.dialogs import status

class REPORTGEN():
    def __init__(self,report_date,
                 report_editor,
                 report_to_customer,
                 product_id,lot_id,
                 stage_id,
                 rawdatas_for_table,
                 nsite,
                 nslot,
                 df2,
                 nset_per_pages):
        self.report_date=report_date
        #print(self.report_date)
        self.report_editor=report_editor
        #print(self.report_editor)
        self.report_to_customer=report_to_customer
        #print(self.report_to_customer)
        self.product_id=product_id
        #print(self.product_id)
        self.lot_id=lot_id
        #print(self.lot_id)
        self.stage_id=stage_id
        #print(self.stage_id)
        self.rawdatas_for_table = rawdatas_for_table
        self.nsite = nsite
        self.nslot = nslot
        ## 從模板檔案 建立 powerpoint物件: => prs
        self.prs=Presentation(" Metrology service report for-[Customer1]_[date].pptx")
        self.df2 = df2
        self.nset_per_page = nset_per_pages
        
    def add_title(self, slide, msg, left, top, width, height):
        """
        args
            slide[slide]: Slide object
            msg[str] : Slide title message
        return:
            None
        """
        shapes = slide.shapes
        shapes.left = Cm(left)
        shapes.top = Cm(top)
        shapes.width = Cm(width)
        shapes.height = Cm(height)
        shapes.title.text = msg
        return shapes

    def add_text(self, slide, msg, left, top, width, height, font_size, is_bold):
        """
        args:
            slide[slide]: Slide object
            msg[str]: Text box message
            left[int]: Position from the left end
            top[int] : Position from top
            width[int]: Width of object
            height[int]: Height of object
            font_size[int]: Font size
            is_bold[int]: Whether to make the text bold
        return:
            None
        """
        textbox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
        #p = textbox.textframe.add_paragraph()
        textbox.text_frame.clear
        p = textbox.text_frame.paragraphs[0]
        p.text = msg
        p.font.size = Pt(font_size)
        p.font.bold = is_bold

        return textbox

    def add_table(self, slide, df, left, top, width, height, font_size):
        '''
        args:
            slide[slide]: Slide object
            df[DataFrame] : Display data
            left[int]: Position from the left end
            top[int] : Position from top
            width[int]: Width of object
            height[int]: Height of object
            font_size[int]: Font size
        return:
            None
        '''

        column_names = df.columns.tolist()
        index_names = df.index.tolist()
        col_num = len(column_names) + 1
        row_num = len(index_names) + 1

        table = slide.shapes.add_table(
             row_num, 
             col_num, 
             Cm(left), 
             Cm(top), 
             Cm(width), 
             Cm(height)
        ).table

        for i in range(col_num):
            table.columns[i].width = Cm(width/col_num)
            if i > 0:
                table.cell(0, i).text = str(column_names[i-1])
                table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(font_size)
                table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for i in range(row_num):
            table.rows[i].height = Cm(height/row_num)
            if i > 0:
                table.cell(i, 0).text = str(index_names[i-1])
                table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(font_size)
                table.cell(i, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for i in range(1, row_num):
            for j in range(1, col_num):
                table.cell(i, j).text = str(df.iloc[i-1, j-1])
                table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(font_size)
                table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return table

    def add_img(self, slide, img_path, left, top, width, height):
        """
        args:
            slide[slide]: Slide object
            img_path[str] : Image file path
            left[int]: Position from the left end
            top[int] : Position from top
            width[int]: Width of object
            height[int]: Height of object 
        return:
            None       
        """
        pic = slide.shapes.add_picture(img_path, 0, 0)
        pic.width = Cm(width)
        pic.height = Cm(height)
        pic.left = Cm(left)
        pic.top = Cm(top)
        return pic

    def add_line(self, slide, bx, by, ex, ey):
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Cm(bx), Cm(by), Cm(ex), Cm(ey))
        return connector

    def go(self):
        output_image_path = "tn"
        ##prs=Presentation("Metrology service report for-[Customer1]_[date].pptx") <=已經移到 REPORTGEN.__init__()中.
        ## ===== A. 封面 =====
        title_slide_layout=self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Technical Support Metrology service "
        subtitle.text = "Date: {} ".format(self.report_date) + "\n" + "FAE: {}".format(self.report_editor) + "\n" + "Customer: {}".format(self.report_to_customer)
        # ===== B. 資訊頁面 =====
        bullet_slide_layout = self.prs.slide_layouts[14]
        slide = self.prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        self.add_title(slide,"Topography Summary Table ",1,0.1,14,1)

        ## 插入產品id, 批號id, 製程id:
        text_info = "Customer product id: {}".format(self.product_id) + "\n" + "Lot id: {}".format(self.lot_id) + "\n" + "Stage: {}".format(self.stage_id)
        self.add_text(slide, text_info,1,3.5,11,3,20,1)

        # print(df)
        ## 插入 summary table: 
        self.add_table(slide, self.df2, 1, 6.5, 10, 6, 14)
        
        # ===== C.把圖檔插入 powerpoint =====
        # 
        # ## 2020/10/28 增加每頁放 nset_per_page組數據.
        #self.nset_per_page = npages
        print("[REPORTGEN]每頁放幾組: ",self.nset_per_page)
        ## 共有幾頁: ex: nslot=13 => int(13/2)=6, 13%2=1 ; 6+1=7頁
        ##              nslot=3 => int(3/2)=1, 3%2=1; 1+1=2pages
        ##     equations: 
        ##                npage = nslot / nset_per_page + nslot%nset_per_page
        self.npage = int(self.nslot / self.nset_per_page) + self.nslot%self.nset_per_page
        print("[REPORTGEN]769: ", "共有幾頁:npage= ", self.npage)
        ## 每頁要放幾組資料: #2020/10/7 修正數據組數計數
        if self.nslot <= 2:
            self.ndata = self.npage * self.nsite * self.nslot
        else:
            self.ndata = self.npage * self.nsite
        print("[REPORTGEN]775: ", "資料數:ndata= ", self.ndata)
        
        ## 判斷量測類型字串: ERO = Erosion , DIS = Dishing , 無定義 = Step Height # 2020/10/28 新增小寫判斷
        if (self.rawdatas_for_table[0]['mtype'] == "ERO") or (self.rawdatas_for_table[0]['mtype'] == "ero"):
            str_mtype = "Erosion "
        elif (self.rawdatas_for_table[0]['mtype'] == "DIS") or (self.rawdatas_for_table[0]['mtype'] == "dis") or (self.rawdatas_for_table[0]['mtype'] == "[J]") or (self.rawdatas_for_table[0]['mtype'] == "[j]"):
            str_mtype = "Dishing "
        else:
            str_mtype ="Step Height "
        
        
        
        
        ##＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊
        #print(" +++++ print out rawdatas_for_table: +++++")
        ## 以 set取得片數的 list: 2021/01/14
        ##   遍例一遍所有資料, 將 slot 放入 list中,
        ##   然後將 list轉為 set, 將重複的 slot整合起來:
        slots_tmp = []
        slots_set = set()
        for ele in self.rawdatas_for_table:
            #print(ele)
            slots_tmp.append(ele['slot']) 
        slots_set = set(slots_tmp)
        print("[REPORTGEN]", slots_set)        
        
        ## raw data 重新整理:
        rawdata_for_report = {}
        rawdatas_for_report = []
        
        for slot_num in slots_set:
            rawdata_for_report['slot'] = slot_num
            datavalues = {}
            datavalue = {}
            for ele in self.rawdatas_for_table:
                if ele['slot'] == slot_num:
                    #print(ele)
                    datavalue['value'] = ele['StepH']
                    datavalue['filename'] = ele['filename']
                    datavalue['mtype'] = ele['mtype']
                    datavalues[ele['site']] = datavalue
                    datavalue={}
            rawdata_for_report['datavalues'] = datavalues
            rawdatas_for_report.append(rawdata_for_report)
            rawdata_for_report={}    
            # print(rawdatas_for_report)
            
            #print('----------')
        #
        #=================================================
        #    rawdatas = [
        #                {  "slot": slot_no ,
        #                   "datavalue": {
        #                                 "C": {"value": step height value,
        #                                       "filename": filename string} ,
        #                                 "M": {"value": step height value,
        #                                       "filename": filename string} ,
        #                                 "E": {"value": step height value,
        #                                       "filename": filename string} ,
        #                                 "1": {"value": step height value,
        #                                       "filename": filename string} ,
        #                                 "2": {"value": step height value,
        #                                       "filename": filename string} ,
        #                                 "3": {"value": step height value,
        #                                       "filename": filename string} ,
        #                                 ...  
        #                                 }  
        #                }
        #               ]
        #
        # for ele in rawdatas_for_report:
        #     print(ele)
        
        ## page = 0 ~ npage-1 ( ie. npage = 2頁, 分別為 page0, page1... 以此類推)
        slides_page=[]
        for page in range(0,self.npage):
            blank_slide_layout = self.prs.slide_layouts[15]
            slides_page.append(self.prs.slides.add_slide(blank_slide_layout))
        ## 迴圈中 dataset := [0,1,2,3,4,5..., rawdatas_for_report list長度 -1 ]
        for dataset in range(0,len(rawdatas_for_report)):
            
            #blank_slide_layout = self.prs.slide_layouts[15]

            #slide = self.prs.slides.add_slide(blank_slide_layout)
            curr_page = int(dataset / self.nset_per_page)
            slide = slides_page[curr_page]
            ## ----------------------------------------------------------------------------------------------------------
            self.add_text(slide, "Topography: {}".format(str_mtype),1,1,6,1,28,1)
            
            space = 0.4
            left = top = 3.5
            width = 6.8
            height = 4.6
            if dataset % self.nset_per_page == 1:
                left = 3.5
                top = 9.5
            ## 迴圈中 ele_site := rawdatas_for_report[第幾組]['datavalues'] -> { site1:{}, site2:{}, site3:{},...}
            ##        ele_site = <str> '1'->'2'->'3'...
            ##        以 ele_site做 iterator直接查詢 rawdatas_for_report字典取值.
            i=0
            for ele_site in rawdatas_for_report[dataset]['datavalues']:
                ## 新增片數標籤:
                #print("Debug#824:", ele_site, type(ele_site))
                lbl=self.add_text(slide, "#{:02d}".format(rawdatas_for_report[dataset]['slot']), 0, top, 1.7, 1.11, 20, 0)
                lbl.fill.patterned()
                lbl.fill.fore_color.rgb=RGBColor(57, 38, 80)
                lbl.fill.back_color.rgb=RGBColor(57, 38, 80)
                lbl.text_frame.paragraphs[0].font.color.rgb=RGBColor(255,255,255)
                ## 按照 site依序插入 image file: eg. 'C' -> 'M' -> 'E' -> 'FE' or '1' -> '2' -> '3' -> '4' -> '5'
                #self.add_img(slide, output_image_path + "\\"+rawdatas_for_report[dataset]['datavalues'][ele_site]['filename'],   left, top, width, height)
                self.add_img(slide, output_image_path + "/"+rawdatas_for_report[dataset]['datavalues'][ele_site]['filename'],   left, top, width, height)
                ## 新增 文字方塊 for site:
                lblsite = self.add_text(slide, "{}".format(ele_site) , left, top+height-1.0, 6, 1, 16, 0)
                lblsite.text_frame.paragraphs[0].font.color.rgb=RGBColor(193,216,47) ## 亮綠色
                ## 新增 文字方塊 for data:
                self.add_text(slide, "{}: {}A".format(str_mtype, rawdatas_for_report[dataset]['datavalues'][ele_site]['value']) , left, top+height, 6, 1, 11, 0)
                ## 新增 線條:
                ### 2021/6/16 新增程式內標線是否顯示之開關變數:
                redline=1 
                blueline=0
                if redline:
                    line1=self.add_line(slide, left+0.5, top+1, left+1.5, top+1)
                    line1.line.color.rgb=RGBColor(255,0,0)
                    line1.line.width=Pt(2)
                    line2=self.add_line(slide, left+1.5, top+3, left+2.5, top+3)
                    line2.line.color.rgb=RGBColor(255,0,0)
                    line2.line.width=Pt(2)

                if blueline:
                    line3=self.add_line(slide, left+1.5, top+1, left+1.5, top+3)
                    line3.line.color.rgb=RGBColor(5,0,255)
                    line3.line.width=Pt(1)
                
                left = left + width + space
                
            i=i+1
            print(i)
            ## ----------------------------------------------------------------------------------------------------------
        print("[REPORTGEN] ＊＊＊＊ Report cooked ＊＊＊＊")
        #print(self.rawdatas_for_table)

    def save_report(self):
        self.prs.save("test-myoutput.pptx")
        print("[REPORTGEN] ===== Report saved ! =====")
##＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊
##＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊＊

#234567890123456789012345678901234567890123456789012345678901234567890123456789012345678901234567890    
# 
# 
# 
# 
# 
# 
    
def main():
    window = tk.Tk()
    window.geometry("640x240+0+0")
    window.title("量測結果報告產生器 v6")
        
    app = Application(master=window)
    app.mainloop()

if __name__ == '__main__':
    main()  # 或是任何你想執行的函式
